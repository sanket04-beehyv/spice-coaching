"""Programmatic test-fixture builders for W-2 Stage A.

Generates synthetic PDFs, PPTXs, and DOCXs in a tmp_path so tests don't have
to commit binary fixtures. The legacy-Bijoy Bangla case is harder to fake
synthetically, so the corresponding test points at the real BRAC SK Basic
Training PDF when available; the test gracefully skips if it's not present.
"""

from __future__ import annotations

from pathlib import Path

import pymupdf  # type: ignore[import-untyped]
from docx import Document  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]

# ── PDF builders ───────────────────────────────────────────────────────


def build_clean_english_pdf(out_path: Path, page_count: int = 3) -> Path:
    """Build a multi-page English PDF with clean Unicode text and headings."""
    doc = pymupdf.open()
    for i in range(page_count):
        page = doc.new_page()
        page.insert_text(
            (72, 72),
            f"# Heading for Page {i + 1}\n"
            f"This is a clean English page about clinical training. "
            f"Page number is {i + 1}.\n"
            f"It contains enough characters to pass the text-empty check.\n"
            f"BP threshold reference: 140/90. Hb reference: < 8 g/dL.",
            fontsize=11,
        )
    doc.save(str(out_path))
    doc.close()
    return out_path


def build_unicode_bangla_pdf(out_path: Path, page_count: int = 3) -> Path:
    """Build a multi-page Bangla PDF using a Unicode-correct font.

    Uses pymupdf's built-in font fallback — text is encoded as UTF-8, so the
    extracted bytes will be valid Bangla codepoints (passes Bangla heuristic).
    """
    doc = pymupdf.open()
    bangla_text = (
        "# অধ্যায় শিরোনাম\n"
        "এটি একটি বাংলা পৃষ্ঠা যেখানে ক্লিনিক্যাল প্রশিক্ষণ সম্পর্কে তথ্য রয়েছে।\n"
        "BP থ্রেশহোল্ড: 140/90। Hb থ্রেশহোল্ড: 8 g/dL এর কম।\n"
        "যথেষ্ট অক্ষর সংখ্যা ক্যারেক্টার এম্পটি চেক পাস করার জন্য।"
    )
    for i in range(page_count):
        page = doc.new_page()
        # Use HtmlBox so unicode bangla characters render properly via pymupdf's
        # built-in font selection.
        page.insert_htmlbox(
            pymupdf.Rect(36, 36, 540, 800),
            f"<p>{bangla_text}</p><p>পৃষ্ঠা নম্বর {i + 1}</p>",
        )
    doc.save(str(out_path))
    doc.close()
    return out_path


def build_image_only_pdf(out_path: Path, page_count: int = 2) -> Path:
    """Build a PDF whose pages contain only rendered images (no extractable text).

    Each page is a blank PDF page with a small inserted Pixmap. pymupdf
    get_text returns near-empty for these pages → text-empty heuristic fails
    → vision fallback in real pipeline.
    """
    doc = pymupdf.open()
    for _ in range(page_count):
        page = doc.new_page()
        # Insert a tiny solid-colour pixmap (no text to extract)
        pix = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 50, 50))
        pix.clear_with(200)
        page.insert_image(pymupdf.Rect(100, 100, 200, 200), pixmap=pix)
    doc.save(str(out_path))
    doc.close()
    return out_path


def build_corrupted_pdf(out_path: Path) -> Path:
    """Write deliberately malformed bytes so pymupdf.open raises."""
    out_path.write_bytes(b"%PDF-1.4\n<<this is not valid pdf content>>\n%%EOF")
    return out_path


def build_empty_pdf(out_path: Path) -> Path:
    """Zero-page PDF — hand-crafted minimal byte stream.

    pymupdf refuses to save a doc with 0 pages, so we emit a minimal valid
    PDF structure with `Pages.Count = 0` directly. pymupdf can open this and
    reports `len(doc) == 0`.
    """
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<</Type/Catalog/Pages 2 0 R>>\nendobj\n"
        b"2 0 obj\n<</Type/Pages/Count 0/Kids[]>>\nendobj\n"
        b"xref\n"
        b"0 3\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000054 00000 n \n"
        b"trailer\n<</Size 3/Root 1 0 R>>\n"
        b"startxref\n100\n%%EOF\n"
    )
    out_path.write_bytes(pdf_bytes)
    return out_path


# ── PPTX builder ───────────────────────────────────────────────────────


def build_clean_english_pptx(out_path: Path, slide_count: int = 3) -> Path:
    """Build a small pptx with title + body text per slide."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for i in range(slide_count):
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = f"UHIS Slide {i + 1}: Clinical Topic"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"This is slide {i + 1} body content."
        tf.add_paragraph().text = "BP reference: 140/90. Hb reference: 8 g/dL."
        tf.add_paragraph().text = "Sufficient body characters for text extraction."
    prs.save(str(out_path))
    return out_path


# ── DOCX builder ───────────────────────────────────────────────────────


def build_clean_english_docx(out_path: Path) -> Path:
    """Build a small docx with Heading 1, Heading 2, and body paragraphs."""
    doc = Document()
    doc.add_heading("BRAC SK Manual: Antenatal Care", level=1)
    doc.add_paragraph("Antenatal care is a set of services for pregnant women.")
    doc.add_heading("Risk Factors", level=2)
    doc.add_paragraph("Common risk factors include high BP and low haemoglobin.")
    doc.add_paragraph("BP threshold: 140/90. Hb threshold: 8 g/dL.")
    doc.add_heading("Referral Decisions", level=2)
    doc.add_paragraph("Refer high-risk PWs to the appropriate facility.")
    doc.save(str(out_path))
    return out_path


def build_formatted_english_docx(out_path: Path) -> Path:
    """DOCX with bold runs and bullet-list paragraphs."""
    doc = Document()
    doc.add_heading("Formatted Manual", level=1)
    para = doc.add_paragraph()
    para.add_run("Normal text and ")
    bold_run = para.add_run("bold emphasis")
    bold_run.bold = True
    doc.add_paragraph("First bullet item", style="List Bullet")
    doc.add_paragraph("Second bullet item", style="List Bullet")
    doc.add_paragraph("Numbered step one", style="List Number")
    doc.add_paragraph("Numbered step two", style="List Number")
    doc.save(str(out_path))
    return out_path


def build_formatted_english_pptx(out_path: Path) -> Path:
    """PPTX slide with bold text in the content placeholder."""
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Formatted Slide"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    opening = tf.paragraphs[0]
    run1 = opening.add_run()
    run1.text = "Opening line with "
    bold = opening.add_run()
    bold.text = "bold lead-in"
    bold.font.bold = True
    bullet = tf.add_paragraph()
    bullet.text = "Bulleted clinical point"
    bullet.level = 1
    prs.save(str(out_path))
    return out_path
