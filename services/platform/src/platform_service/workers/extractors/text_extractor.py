"""W-2 Stage A — text-first extraction.

Per Pipeline v3.3 §4.1. Source-type dispatcher:
- PDF → pymupdf get_text per page
- PPTX → python-pptx slide text per slide (one slide = one "page")
- DOCX → python-docx paragraph text (one document = one "page" — docx
  doesn't have inherent pagination; reviewers drill into the whole doc)

Returns markdown-ish content with light heading hierarchy where the source
exposes it (PPTX slide title, DOCX heading style, PDF text font hints).

The output of this module feeds the quality heuristic, which decides whether
vision fallback is needed.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

import pymupdf  # type: ignore[import-untyped]
from docx import Document as DocxDocument  # type: ignore[import-untyped]
from docx.text.paragraph import Paragraph as DocxParagraph  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]

from platform_service.workers.extractors.extraction_markdown import normalize_extraction_markdown

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class ExtractedPage:
    """Per-page text or transcript extraction result.

    Document extractors (PDF/Office) populate only ``page_number`` and
    ``markdown``. Transcript extractors additionally populate ``start_ms``
    / ``end_ms`` (deterministic chunk boundaries from the AV splitter) and
    may populate ``extraction_quality_score`` (a composite measure of the
    transcript's coherence) and ``language_detected``. The orchestrator
    persists these fields to ``source_page`` when present.
    """

    page_number: int  # 1-indexed
    markdown: str
    start_ms: int | None = None
    end_ms: int | None = None
    extraction_quality_score: float | None = None
    language_detected: str | None = None


class TextExtractionError(Exception):
    """Raised when text extraction fails irrecoverably (corrupted file, etc.)."""


def _finalize_markdown(markdown: str) -> str:
    return normalize_extraction_markdown(markdown)


def _docx_run_text(run) -> str:
    text = run.text or ""
    if not text:
        return ""
    if run.bold:
        text = f"**{text}**"
    if run.italic:
        text = f"*{text}*"
    return text


def _docx_paragraph_text(para: DocxParagraph) -> str:
    parts = [_docx_run_text(run) for run in para.runs]
    combined = "".join(parts).strip()
    if combined:
        return combined
    return (para.text or "").strip()


def _docx_list_prefix(para: DocxParagraph, *, num_counters: dict[int, int]) -> str | None:
    style_name = (para.style.name or "").strip() if para.style else ""
    if style_name.startswith("List Bullet"):
        return "- "
    if style_name.startswith("List Number"):
        num_counters[0] = num_counters.get(0, 0) + 1
        return f"{num_counters[0]}. "
    p_pr = para._p.pPr
    if p_pr is not None and p_pr.numPr is not None:
        return "- "
    return None


def _docx_paragraph_line(para: DocxParagraph, *, num_counters: dict[int, int]) -> str:
    text = _docx_paragraph_text(para)
    if not text:
        return ""
    prefix = _docx_list_prefix(para, num_counters=num_counters)
    if prefix is not None:
        return f"{prefix}{text}"
    num_counters.clear()
    return text


def _pptx_run_text(run) -> str:
    text = run.text or ""
    if not text:
        return ""
    if run.font.bold:
        text = f"**{text}**"
    if run.font.italic:
        text = f"*{text}*"
    return text


def _pptx_paragraph_text(paragraph) -> str:
    parts = [_pptx_run_text(run) for run in paragraph.runs]
    combined = "".join(parts).strip()
    if combined:
        return combined
    return (paragraph.text or "").strip()


def _pptx_is_bulleted(paragraph) -> bool:
    if (paragraph.level or 0) > 0:
        return True
    p_pr = paragraph._p.pPr
    if p_pr is None:
        return False
    ns = "{http://schemas.openformats.org/drawingml/2006/main}"
    return any(p_pr.find(f"{ns}{tag}") is not None for tag in ("buChar", "buAutoNum", "buBlip"))


def _pptx_paragraph_line(paragraph) -> str:
    text = _pptx_paragraph_text(paragraph)
    if not text:
        return ""
    level = paragraph.level or 0
    indent = "  " * level
    if _pptx_is_bulleted(paragraph):
        return f"{indent}- {text}"
    return text


# ── PDF ────────────────────────────────────────────────────────────────


def extract_pdf_pages(pdf_path: str | Path) -> list[ExtractedPage]:
    """Extract markdown-ish text from each PDF page via pymupdf.

    Empty PDFs return []. Corrupted PDFs raise TextExtractionError.
    """
    path = Path(pdf_path)
    if not path.is_file():
        raise TextExtractionError(f"PDF not found: {path}")
    try:
        doc = pymupdf.open(str(path))
    except Exception as exc:
        raise TextExtractionError(f"Failed to open PDF {path.name}: {exc}") from exc
    try:
        pages = []
        for i, page in enumerate(doc):
            try:
                text = page.get_text() or ""
            except Exception:
                # Single bad page should not abort the whole document.
                logger.warning(
                    "pymupdf get_text failed for page %d of %s — using empty",
                    i + 1,
                    path.name,
                )
                text = ""
            pages.append(ExtractedPage(page_number=i + 1, markdown=_finalize_markdown(text.strip())))
        return pages
    finally:
        doc.close()


# ── PPTX ───────────────────────────────────────────────────────────────


def extract_pptx_slides(pptx_path: str | Path) -> list[ExtractedPage]:
    """Extract text from each pptx slide as one "page".

    Slide title becomes a markdown # heading; non-title text frames become
    paragraphs. Embedded videos / audio / pictures are skipped (no text).
    Master slide content is excluded (per Implementation Plan W-2 edge case 7).
    """
    path = Path(pptx_path)
    if not path.is_file():
        raise TextExtractionError(f"PPTX not found: {path}")
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise TextExtractionError(f"Failed to open PPTX {path.name}: {exc}") from exc

    pages: list[ExtractedPage] = []
    for slide_idx, slide in enumerate(prs.slides):
        title_text = ""
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            shape_lines = [line for p in tf.paragraphs if (line := _pptx_paragraph_line(p))]
            shape_text = "\n".join(shape_lines)
            if not shape_text:
                continue
            # Title placeholder identification: slides title style is at index
            # 0 typically; we approximate by shape.is_placeholder + idx 0.
            is_title = bool(
                getattr(shape, "is_placeholder", False)
                and getattr(shape.placeholder_format, "idx", None) == 0
            )
            if is_title and not title_text:
                title_text = shape_text
            else:
                body_lines.append(shape_text)
        # Build markdown
        parts: list[str] = []
        if title_text:
            parts.append(f"# {title_text}")
        if body_lines:
            parts.extend(body_lines)
        markdown = _finalize_markdown("\n\n".join(parts).strip())
        pages.append(ExtractedPage(page_number=slide_idx + 1, markdown=markdown))
    return pages


# ── DOCX ───────────────────────────────────────────────────────────────


def extract_docx_pages(docx_path: str | Path) -> list[ExtractedPage]:
    """Extract docx as a single logical "page" with markdown-ish heading levels.

    DOCX has no inherent pagination; for ingestion we treat the whole
    document as page 1. Heading paragraphs become #/##/### markdown headings;
    other paragraphs become plain lines.
    """
    path = Path(docx_path)
    if not path.is_file():
        raise TextExtractionError(f"DOCX not found: {path}")
    try:
        doc = DocxDocument(str(path))
    except Exception as exc:
        raise TextExtractionError(f"Failed to open DOCX {path.name}: {exc}") from exc

    lines: list[str] = []
    num_counters: dict[int, int] = {}
    for para in doc.paragraphs:
        text = _docx_paragraph_line(para, num_counters=num_counters)
        if not text:
            continue
        style_name = (para.style.name or "").strip() if para.style else ""
        if style_name.startswith("Heading"):
            num_counters.clear()
            # Heading 1 → "# ", Heading 2 → "## ", etc. Default to ## for any
            # malformed heading style.
            try:
                level = int(style_name.split()[-1])
                level = max(1, min(level, 6))
            except (ValueError, IndexError):
                level = 2
            heading_text = _docx_paragraph_text(para)
            lines.append(f"{'#' * level} {heading_text}")
        else:
            lines.append(text)
    markdown = _finalize_markdown("\n\n".join(lines).strip())
    return [ExtractedPage(page_number=1, markdown=markdown)]


# ── Dispatcher ─────────────────────────────────────────────────────────


def extract_pages(source_path: str | Path, source_type: str) -> list[ExtractedPage]:
    """Source-type-aware text extraction. Returns one ExtractedPage per
    logical page (slide for pptx, document for docx)."""
    source_type = source_type.lower()
    if source_type == "pdf":
        return extract_pdf_pages(source_path)
    if source_type == "pptx":
        return extract_pptx_slides(source_path)
    if source_type == "docx":
        return extract_docx_pages(source_path)
    raise TextExtractionError(f"Unsupported source_type: {source_type!r}")
